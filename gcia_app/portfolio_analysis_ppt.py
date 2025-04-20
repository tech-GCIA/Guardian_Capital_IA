from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from django.contrib.staticfiles.finders import find
from pptx.enum.shapes import MSO_SHAPE
from gcia_app.utils import format_currency, calculate_portfolio_and_benchmark_xirr

def add_custom_header(slide, title_text, subtitle_text, logo_path, summary_slide=False):
    """Adds a header to the slide with a title, subtitle, black line, and logo.
    Modified for the updated slide width."""
    # Title (Green)
    title_left = Inches(0.5)
    title_top = Inches(0.1)
    title_width = Inches(12)  # Increased width for the title
    title_height = Inches(0.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.add_paragraph()
    title_para.text = title_text
    title_para.font.bold = True
    title_para.font.size = Pt(24)
    title_para.font.color.rgb = RGBColor(34, 139, 34)  # Forest Green
    title_para.alignment = PP_ALIGN.LEFT

    # Subtitle (Italic Blue)
    subtitle_top = title_top + Inches(0.5)
    subtitle_height = Inches(0.3)
    subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_para = subtitle_frame.add_paragraph()
    subtitle_para.text = subtitle_text
    subtitle_para.font.italic = True
    subtitle_para.font.size = Pt(12)
    subtitle_para.font.color.rgb = RGBColor(0, 0, 128)  # Navy Blue
    subtitle_para.alignment = PP_ALIGN.LEFT

    # Black Line Separator - adjusted for wider slide
    line_top = subtitle_top + Inches(0.3)
    line_width = Inches(14)  # Increased width for the line
    line_height = Pt(1.5)  # Thin black line
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), line_top, line_width, line_height
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color
    line_shape.line.fill.background()  # Remove border

    # Logo (Aligned with title bottom) - adjusted position for wider slide
    logo_width = Inches(1.5)
    logo_left = Inches(13)  # Adjusted position for wider slide
    logo_top = Inches(0.2)
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)

    if not summary_slide:
        # Labels (Rounded Rectangles)
        label_texts = ["1. Concentration", "2. Quality", "3. Price"]
        label_top = line_top + Inches(0.3)
        label_width = Inches(1.5)
        label_height = Inches(0.25)
        label_spacing = Inches(0.8)  # Increased spacing between labels

        for i, text in enumerate(label_texts):
            label_left = Inches(1.5) + i * (label_width + label_spacing)  # Adjusted starting position
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, label_left, label_top, label_width, label_height
            )
            shape.text = text

            # Format the label
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(200, 230, 200)  # Light green
            line = shape.line
            line.color.rgb = RGBColor(0, 128, 0)  # Dark green
            line.width = Pt(1.5)

            # Center-align the text
            text_frame = shape.text_frame
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)# Black

def create_table_on_slide(slide, data, left, top, width, height, header_color, row_color):
    """Helper function to create a formatted table on a slide"""
    rows = len(data)
    cols = len(data[0])
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    for i, col in enumerate(table.columns):
        if i == 0:
            col.width = Inches(width.inches * 0.4)  # Make first column a bit wider
        else:
            col.width = Inches(width.inches * (0.6 / (cols - 1)))  # Distribute remaining space
    
    # Add headers and format them
    for col_idx, text in enumerate(data[0]):
        cell = table.cell(0, col_idx)
        cell.text = text
        
        # Header formatting
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(11)
        paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Set text color to white
        run = paragraph.runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Set header background color
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
    
    # Add data and format rows
    for row_idx in range(1, rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            
            # Cell formatting
            paragraph = cell.text_frame.paragraphs[0]
            # Adjust font size if we have many rows
            if rows > 15:
                paragraph.font.size = Pt(9)  # Smaller font for many rows
            else:
                paragraph.font.size = Pt(10)
            
            if col_idx == 0:
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Set alternate row colors
            cell.fill.solid()
            if row_idx % 2 == 1:
                cell.fill.fore_color.rgb = row_color
            else:
                cell.fill.fore_color.rgb = RGBColor(235, 240, 247)  # Lighter blue
    
    # Adjust row heights based on number of rows
    row_height = min(height.inches / rows, 0.35)  # Cap at 0.35 inches, but reduce if needed
    for row in table.rows:
        row.height = Inches(row_height)
    
    return table

def add_equity_review_slide(prs, summary_data, logo_path):
    """Add a slide reviewing equity mutual funds with updated styling and width adjustments
    
    Args:
        prs: PowerPoint presentation object
        summary_data: List of dictionaries containing fund details
        logo_path: Path to the logo image
    """
    # Filter for equity funds only
    equity_funds = [fund for fund in summary_data 
                   if fund["Fund Category"].lower().startswith("equity")]
    
    # Determine how many funds we can fit on one slide
    MAX_FUNDS_PER_SLIDE = 12  # Maximum number of funds per slide
    
    # Calculate number of slides needed
    num_slides = (len(equity_funds) + MAX_FUNDS_PER_SLIDE - 1) // MAX_FUNDS_PER_SLIDE
    
    # Get unique fund categories and count
    fund_categories = {}
    for fund in equity_funds:
        category = fund["Fund Category"]
        if category not in fund_categories:
            fund_categories[category] = 0
        fund_categories[category] += 1
    
    # Process equity funds in batches for each slide
    for slide_idx in range(num_slides):
        # Create new slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add header
        slide_title = "Equity MF Review"
        if num_slides > 1:
            slide_title += f" ({slide_idx + 1}/{num_slides})"
            
        add_custom_header(
            slide, 
            slide_title, 
            "Following is the review of the equity portfolio", 
            logo_path,
            summary_slide=True
        )
        
        # Calculate the start and end indices for this slide
        start_idx = slide_idx * MAX_FUNDS_PER_SLIDE
        end_idx = min(start_idx + MAX_FUNDS_PER_SLIDE, len(equity_funds))
        
        # Get the funds for this slide
        current_funds = equity_funds[start_idx:end_idx]
        
        # Create table data with headers
        table_data = [["No", "Particulars", "Concentration", "Quality", "Price"]]
        
        # Add fund rows
        for i, fund in enumerate(current_funds):
            absolute_idx = start_idx + i + 1  # For continuous numbering across slides
            table_data.append([
                str(absolute_idx),
                fund["Scheme"],
                fund.get("Concentration", ""),
                fund.get("Quality", ""),
                fund.get("Price", "")
            ])
        
        # Add total row if this is the last slide
        if slide_idx == num_slides - 1:
            table_data.append(["", "Total", "", "", ""])
            table_data.append(["", "Nifty 500", "", "", ""])
        
        # Calculate table dimensions - adjust for new slide width
        table_left = Inches(0.5)
        table_top = Inches(1.35)
        table_width = Inches(14)  # Increased width to match slide width
        
        # Adjust table height to leave room for notes
        if slide_idx == num_slides - 1:  # Last slide needs space for notes
            table_height = Inches(3.2)   # Shorter table on the last slide to fit notes
        else:
            table_height = Inches(4.5)   # Full height table on other slides
        
        # Create the table
        header_color = RGBColor(65, 105, 225)  # Royal blue for headers
        row_color = RGBColor(220, 230, 242)    # Light blue for rows
        
        table = create_table_on_slide(
            slide, 
            table_data, 
            table_left, 
            table_top, 
            table_width, 
            table_height,
            header_color,
            row_color
        )
        
        # Set column widths - adjust as needed for wider slide
        if table:
            # Wider column for fund names
            table.columns[0].width = Inches(0.7)  # No column
            table.columns[1].width = Inches(6.8)  # Increased width for Particulars column
            table.columns[2].width = Inches(2.2)  # Slightly increased width for Concentration
            table.columns[3].width = Inches(2.2)  # Slightly increased width for Quality
            table.columns[4].width = Inches(2.1)  # Slightly increased width for Price
        
        # Add notes section on the last slide - adjust width
        if slide_idx == num_slides - 1:
            notes_top = Inches(5.6)  # Fixed position instead of relative to table
            notes_left = Inches(0.5)
            
            # Add Notes heading
            notes_title = slide.shapes.add_textbox(
                left=notes_left, 
                top=notes_top, 
                width=Inches(14),  # Increased width
                height=Inches(0.3)
            )
            notes_title.text_frame.text = "Notes:"
            notes_title.text_frame.paragraphs[0].font.bold = True
            
            # Add bullet points
            bullet_points = [
                f"Marketcap Allocation",
                f"No of funds in each category",
                f"No of funds concentrated",
                f"No of funds over/under valued",
                f"No of funds performing",
                f"Ideal no of funds and total stocks"
            ]
            
            bullets_box = slide.shapes.add_textbox(
                left=notes_left, 
                top=Inches(5.9), 
                width=Inches(14),  # Increased width
                height=Inches(1.5)
            )
            
            # Add each bullet point
            text_frame = bullets_box.text_frame
            for i, point in enumerate(bullet_points):
                if i == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()
                
                paragraph.text = point
                paragraph.level = 0
                paragraph.font.size = Pt(12)
                
                # Add bullet character
                paragraph.font.size = Pt(12)
                paragraph.bullet = True
    return

def add_hybrid_review_slide(prs, summary_data, logo_path):
    """Add a slide reviewing hybrid mutual funds with updated styling and width adjustments
    
    Args:
        prs: PowerPoint presentation object
        summary_data: List of dictionaries containing fund details
        logo_path: Path to the logo image
    """
    # Filter for hybrid funds only
    hybrid_funds = [fund for fund in summary_data 
                   if fund["Fund Category"].lower().startswith("hybrid")]
    
    # Determine how many funds we can fit on one slide
    MAX_FUNDS_PER_SLIDE = 12  # Maximum number of funds per slide
    
    # Calculate number of slides needed
    num_slides = (len(hybrid_funds) + MAX_FUNDS_PER_SLIDE - 1) // MAX_FUNDS_PER_SLIDE
    
    # Process hybrid funds in batches for each slide
    for slide_idx in range(num_slides):
        # Create new slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add header
        slide_title = "Hybrid MF Review"
        if num_slides > 1:
            slide_title += f" ({slide_idx + 1}/{num_slides})"
            
        add_custom_header(
            slide, 
            slide_title, 
            "Following is the review of the hybrid/debt mf portfolio", 
            logo_path,
            summary_slide=True
        )
        
        # Calculate the start and end indices for this slide
        start_idx = slide_idx * MAX_FUNDS_PER_SLIDE
        end_idx = min(start_idx + MAX_FUNDS_PER_SLIDE, len(hybrid_funds))
        
        # Get the funds for this slide
        current_funds = hybrid_funds[start_idx:end_idx]
        
        # Create table data with headers
        table_data = [["No", "Particulars", "No of securities", "Age of the fund", 
                      "Equity Allocation", "Debt Allocation", "Cash Allocation", "Expense ratio", "Credit rating AAA", "Credit rating AA", "Credit rating A", 
                      "Performance", "Up Capture", "Down Capture"]]
        
        # Add fund rows
        for i, fund in enumerate(current_funds):
            absolute_idx = start_idx + i + 1  # For continuous numbering across slides
            
            table_data.append([
                str(absolute_idx),
                fund["Scheme"],
                str(fund.get("No of Securities", "")),
                f"{fund.get('Age of the fund', '')} yrs" if fund.get('Age of the fund', '') else "-",
                f"{fund.get('Equity Allocation', '')}%" if fund.get('Equity Allocation', '') else "-",
                f"{fund.get('Debt Allocation', '')}%" if fund.get('Debt Allocation', '') else "-",
                f"{fund.get('Cash Allocation', '')}%" if fund.get('Cash Allocation', '') else "-",
                f"{fund.get('Expense ratio', '')}%" if fund.get('Expense ratio', '') else "-",
                f"{fund.get('Credit rating AAA', '')}%" if fund.get('Credit rating AAA', '') else "-",
                f"{fund.get('Credit rating AA', '')}%" if fund.get('Credit rating AA', '') else "-",
                f"{fund.get('Credit rating A', '')}%" if fund.get('Credit rating A', '') else "-",
                fund.get("Performance", ""),
                f"{fund.get('Up Capture', '')}%" if fund.get('Up Capture', '') else "-",
                f"{fund.get('Down Capture', '')}%" if fund.get('Down Capture', '') else "-"
            ])
        
        # Add total row if this is the last slide
        if slide_idx == num_slides - 1:
            table_data.append(["", "Total"] + [""] * 12)
            table_data.append(["", "Nifty 50 & Liquid Fund (50:50)"] + [""] * 12)
        
        # Calculate table dimensions - adjust for new slide width
        table_left = Inches(0.5)
        table_top = Inches(1.35)
        table_width = Inches(14)  # Increased width to match slide width
        
        # Adjust table height to leave room for notes
        if slide_idx == num_slides - 1:  # Last slide needs space for notes
            table_height = Inches(3.2)   # Shorter table on the last slide to fit notes
        else:
            table_height = Inches(4.5)   # Full height table on other slides
        
        # Create the table
        header_color = RGBColor(65, 105, 225)  # Royal blue for headers
        row_color = RGBColor(220, 230, 242)    # Light blue for rows
        
        table = create_table_on_slide(
            slide, 
            table_data, 
            table_left, 
            table_top, 
            table_width, 
            table_height,
            header_color,
            row_color
        )
        
        # Set column widths - adjust as needed for wider slide
        if table:
            # Calculate total width excluding first two columns (14 columns total)
            remaining_width = 10.0  # Out of 14 total inches
            
            # Fixed widths for first two columns
            table.columns[0].width = Inches(0.5)   # No column
            table.columns[1].width = Inches(3.5)   # Particulars column
            
            # Distribute remaining width among other columns
            column_count = len(table_data[0]) - 2
            column_width = remaining_width / column_count
            
            for i in range(2, len(table_data[0])):
                table.columns[i].width = Inches(column_width)
        
        # Add notes section on the last slide - adjust width
        if slide_idx == num_slides - 1:
            notes_top = Inches(5.8)  # Fixed position instead of relative to table
            notes_left = Inches(0.5)
            
            # Add Notes heading
            notes_title = slide.shapes.add_textbox(
                left=notes_left, 
                top=notes_top, 
                width=Inches(14),  # Increased width
                height=Inches(0.3)
            )
            notes_title.text_frame.text = "Notes:"
            notes_title.text_frame.paragraphs[0].font.bold = True
            
            # Add bullet points
            bullet_points = [
                f"Split between equity/debt & arbitrage",
                f"No of funds which are 5 years and old",
                f"No of funds which have beaten the index",
                f"Capture ratios"
            ]
            
            bullets_box = slide.shapes.add_textbox(
                left=notes_left, 
                top=Inches(6.1), 
                width=Inches(14),  # Increased width
                height=Inches(1.5)
            )
            
            # Add each bullet point
            text_frame = bullets_box.text_frame
            for i, point in enumerate(bullet_points):
                if i == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()
                
                paragraph.text = point
                paragraph.level = 0
                paragraph.font.size = Pt(12)
                
                # Add bullet character
                paragraph.bullet = True
    return

def add_debt_review_slide(prs, summary_data, logo_path):
    """Add a slide reviewing debt mutual funds with updated styling and width adjustments
    
    Args:
        prs: PowerPoint presentation object
        summary_data: List of dictionaries containing fund details
        logo_path: Path to the logo image
    """
    # Filter for debt funds only
    debt_funds = [fund for fund in summary_data 
                 if fund["Fund Category"].lower().startswith("debt")]
    
    # Determine how many funds we can fit on one slide
    MAX_FUNDS_PER_SLIDE = 12  # Maximum number of funds per slide
    
    # Calculate number of slides needed
    num_slides = (len(debt_funds) + MAX_FUNDS_PER_SLIDE - 1) // MAX_FUNDS_PER_SLIDE
    
    
    # Get unique fund categories and count
    fund_categories = {}
    for fund in debt_funds:
        category = fund["Fund Category"]
        if category not in fund_categories:
            fund_categories[category] = 0
        fund_categories[category] += 1
    
    # Process debt funds in batches for each slide
    for slide_idx in range(num_slides):
        # Create new slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add header
        slide_title = "Debt MF Review"
        if num_slides > 1:
            slide_title += f" ({slide_idx + 1}/{num_slides})"
            
        add_custom_header(
            slide, 
            slide_title, 
            "Following is the review of the debt mf portfolio", 
            logo_path,
            summary_slide=True
        )
        
        # Calculate the start and end indices for this slide
        start_idx = slide_idx * MAX_FUNDS_PER_SLIDE
        end_idx = min(start_idx + MAX_FUNDS_PER_SLIDE, len(debt_funds))
        
        # Get the funds for this slide
        current_funds = debt_funds[start_idx:end_idx]
        
        # Create table data with headers
        table_data = [["No", "Particulars", "No of securities", "Mod duration", "YTM(%)", 
                       "Expense ratio", "AAA", "AA", "A"]]
        
        # Add fund rows
        for i, fund in enumerate(current_funds):
            absolute_idx = start_idx + i + 1  # For continuous numbering across slides
            
            table_data.append([
                str(absolute_idx),
                fund["Scheme"],
                str(fund.get("No of Securities", "")),
                str(fund.get("Mod duration", "")),
                f"{fund.get('YTM', '')}%" if fund.get('YTM', '') else "-",
                f"{fund.get('Expense ratio', '')}%" if fund.get('Expense ratio', '') else "-",
                f"{fund.get('Credit rating AAA', '')}%" if fund.get('Credit rating AAA', '') else "-",
                f"{fund.get('Credit rating AA', '')}%" if fund.get('Credit rating AA', '') else "-",
                f"{fund.get('Credit rating A', '')}%" if fund.get('Credit rating A', '') else "-"
            ])
        
        # Add total row if this is the last slide
        if slide_idx == num_slides - 1:
            table_data.append(["", "Total"] + [""] * 7)
        
        # Calculate table dimensions - adjust for new slide width
        table_left = Inches(0.5)
        table_top = Inches(1.35)
        table_width = Inches(14)  # Increased width to match slide width
        
        # Adjust table height to leave room for notes
        if slide_idx == num_slides - 1:  # Last slide needs space for notes
            table_height = Inches(3.2)   # Shorter table on the last slide to fit notes
        else:
            table_height = Inches(4.5)   # Full height table on other slides
        
        # Create the table
        header_color = RGBColor(65, 105, 225)  # Royal blue for headers
        row_color = RGBColor(220, 230, 242)    # Light blue for rows
        
        table = create_table_on_slide(
            slide, 
            table_data, 
            table_left, 
            table_top, 
            table_width, 
            table_height,
            header_color,
            row_color
        )
        
        # Set column widths - adjust as needed for wider slide
        if table:
            # Calculate total width excluding first two columns (9 columns total)
            remaining_width = 8.5  # Out of 14 total inches
            
            # Fixed widths for first two columns
            table.columns[0].width = Inches(0.5)   # No column
            table.columns[1].width = Inches(5)     # Particulars column (wider for fund names)
            
            # Distribute remaining width among other columns
            column_count = len(table_data[0]) - 2
            column_width = remaining_width / column_count
            
            for i in range(2, len(table_data[0])):
                table.columns[i].width = Inches(column_width)
        
        # Add notes section on the last slide - adjust width
        if slide_idx == num_slides - 1:
            notes_top = Inches(5.6)  # Fixed position instead of relative to table
            notes_left = Inches(0.5)
            
            # Add Notes heading
            notes_title = slide.shapes.add_textbox(
                left=notes_left, 
                top=notes_top, 
                width=Inches(14),  # Increased width
                height=Inches(0.3)
            )
            notes_title.text_frame.text = "Notes:"
            notes_title.text_frame.paragraphs[0].font.bold = True
            
            # Add bullet points
            bullet_points = [
                f"No of funds in each category",
                f"No of funds diversified",
                f"Avg Credit rating",
                f"Total duration of portfolio",
                f"What we can do better"
            ]
            
            bullets_box = slide.shapes.add_textbox(
                left=notes_left, 
                top=Inches(5.9), 
                width=Inches(14),  # Increased width
                height=Inches(1.5)
            )
            
            # Add each bullet point
            text_frame = bullets_box.text_frame
            for i, point in enumerate(bullet_points):
                if i == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()
                
                paragraph.text = point
                paragraph.level = 0
                paragraph.font.size = Pt(12)
                
                # Add bullet character
                paragraph.bullet = True
    return

def add_portfolio_summary_slide(prs, client_details, summary_data, logo_path, transaction_data):
    """Add a summary slide with portfolio details, asset allocation, and fund types
    with updated styling and width adjustments"""
    # Create summary slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add header with logo - adjust width to match slide width
    add_custom_header(
        summary_slide, 
        "Overall Portfolio Summary", 
        "This Slide Review the Mutual Fund Portfolio", 
        logo_path,
        summary_slide=True
    )
    
    # Define table colors
    header_color = RGBColor(65, 105, 225)  # Royal blue for headers
    row_color = RGBColor(220, 230, 242)    # Light blue for rows
    print("transaction_data", transaction_data)

    portfolio_xirr, benchmark_xirr = calculate_portfolio_and_benchmark_xirr(transaction_data)
    # Create left table for Portfolio Particulars - adjust position and size
    left_table_data = [
        ["Particulars", "Value"],
        ["Total Invested", format_currency(round(client_details["Invested Amount"], 1))],
        ["Current Value", format_currency(round(client_details["Current Value"],1))],
        ["Holding period", f"{client_details['Holding Days']} days"],
        ["Absolute Gains", format_currency(round(client_details["Gain"], 1))],
        ["Abs[%]", f"{client_details['Absolute Return']:.1f}%"],
        ["CAGR[%]", f"{client_details['CAGR(%)']:.1f}%"],
        ["XIRR", f"{portfolio_xirr:.1f}%"],
        ["Benchmark XIRR", f"{benchmark_xirr:.1f}%"],
        ["Over/Under performance", "Overperformed" if portfolio_xirr > benchmark_xirr else "Underperformed"]
    ]
    
    # Adjust table position and size to match slide width
    left_table_left = Inches(0.5)
    left_table_top = Inches(1.35)
    left_table_width = Inches(6)  # Increased width
    left_table_height = Inches(3)
    
    create_table_on_slide(
        summary_slide, 
        left_table_data, 
        left_table_left, 
        left_table_top, 
        left_table_width, 
        left_table_height,
        header_color,
        row_color
    )
    
    # Create Type of Funds table
    # Process summary_data to calculate fund types
    fund_types = {"Direct": 0, "Regular": 0}
    
    for fund in summary_data:
        if fund["Is Direct Fund"]:
            fund_types["Direct"] += fund["Current Value"]
        else:
            fund_types["Regular"] += fund["Current Value"]
    
    total_value = client_details["Current Value"]
    
    type_table_data = [
        ["Type of Funds", "Amount", "Allocation", "Funds"],
        ["Direct", format_currency(round(fund_types["Direct"],1)), f"{fund_types['Direct']/total_value*100:.1f}%", 
         str(sum(1 for fund in summary_data if fund["Is Direct Fund"]))],
        ["Regular", format_currency(round(fund_types["Regular"],1)), f"{fund_types['Regular']/total_value*100:.1f}%", 
         str(sum(1 for fund in summary_data if not fund["Is Direct Fund"]))],
        ["Total", format_currency(round(total_value,1)), "100%", str(len(summary_data))]
    ]
    
    # Adjust position to add spacing between tables
    type_table_left = Inches(0.5)
    type_table_top = Inches(4.6)  # Increased spacing from previous table
    type_table_width = Inches(6)  # Increased width to match first table
    type_table_height = Inches(1.5)
    
    create_table_on_slide(
        summary_slide, 
        type_table_data, 
        type_table_left, 
        type_table_top, 
        type_table_width, 
        type_table_height,
        header_color,
        row_color
    )
    
    # Create Asset Allocation table using categories present in summary_data
    # Calculate asset allocation from summary_data
    asset_allocation = {}
    category_fund_counts = {}
    
    # Use summary_data for asset allocation
    for fund in summary_data:
        category = fund["Fund Category"]
        if category not in asset_allocation:
            asset_allocation[category] = 0
            category_fund_counts[category] = 0
        asset_allocation[category] += fund["Current Value"]
        category_fund_counts[category] += 1
    
    # Create category data with percentages for sorting
    category_data = []
    for category, value in asset_allocation.items():
        percentage = (value / total_value * 100) if total_value > 0 else 0
        category_data.append({
            'category': category,
            'value': value,
            'percentage': percentage,
            'fund_count': category_fund_counts[category]
        })
    
    # Sort categories by allocation percentage in descending order
    sorted_categories = sorted(category_data, key=lambda x: x['percentage'], reverse=True)
    
    asset_table_data = [["Asset Allocation", "Amount", "Allocation", "No of funds"]]
    
    # Handle case with many categories
    MAX_ROWS_PER_SLIDE = 15  # Maximum number of category rows to display
    
    # If we have too many categories, we'll take the top ones by allocation
    if len(sorted_categories) > MAX_ROWS_PER_SLIDE:
        # Take top categories up to MAX_ROWS_PER_SLIDE - 1 (leaving room for "Others" row)
        displayed_categories = sorted_categories[:MAX_ROWS_PER_SLIDE - 1]
        other_categories = sorted_categories[MAX_ROWS_PER_SLIDE - 1:]
        
        # Calculate totals for "Others" row
        others_value = sum(cat['value'] for cat in other_categories)
        others_percentage = sum(cat['percentage'] for cat in other_categories)
        others_fund_count = sum(cat['fund_count'] for cat in other_categories)
        
        # Add each top category row
        for cat_data in displayed_categories:
            asset_table_data.append([
                cat_data['category'], 
                format_currency(round(cat_data['value'],1)), 
                f"{cat_data['percentage']:.1f}%",
                str(cat_data['fund_count'])
            ])
        
        # Add "Others" row
        asset_table_data.append([
            "Others", 
            format_currency(round(others_value, 1)), 
            f"{others_percentage:.1f}%",
            str(others_fund_count)
        ])
    else:
        # Add each category row if we have a manageable number
        for cat_data in sorted_categories:
            asset_table_data.append([
                cat_data['category'], 
                format_currency(round(cat_data['value'],1)), 
                f"{cat_data['percentage']:.1f}%",
                str(cat_data['fund_count'])
            ])
    
    
    # Add total row
    asset_table_data.append([
        "Total", 
        format_currency(round(total_value,1)), 
        "100%",
        str(len(summary_data))
    ])
    
    # Adjust position and size for asset allocation table
    asset_table_left = Inches(7)  # Increased spacing between tables
    asset_table_top = Inches(1.35)
    asset_table_width = Inches(7.5)  # Adjusted width to fit slide width
    asset_table_height = Inches(4.85)  # Slightly increased height
    
    create_table_on_slide(
        summary_slide, 
        asset_table_data, 
        asset_table_left, 
        asset_table_top, 
        asset_table_width, 
        asset_table_height,
        header_color,
        row_color
    )
    
    return summary_slide


def create_fund_presentation(ppt_data_dict, client_details, summary_data, transaction_data):
    """format for ppt_data_dict:
        [{
            "Scheme Name": [
                "Axis Blue Chip Fund Reg (G)",
                "BSE 100 TRI ",
                "Equity Largecap"
            ],
            "Purchase Value": 10000,
            "Current Value": 20000,
            "Gain": 10000,
            "Weight": 23,
            "Stocks": 45,
            "YTD": [
                23.32, 25.56, 27.32
            ],
            "1YR": [
                23.32, 25.56, 27.32
            ],
            "3YR": [
                23.32, 25.56, 27.32
            ],
            "5YR": [
                23.32, 25.56, 27.32
            ],
            "PE": 30,
            "Index PE": 40,
            "Comments": [
                "Concentration: Concentrated",
                "Performance: Under performed",
                "Valuation: Overvalued"
            ]
        }]
    """
    prs = Presentation()
    
    # Set slide dimensions
    prs.slide_width = Inches(15)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background image to title slide
    left = top = 0
    image_path = find('gcia_app/images/ppt_image.png')
    logo_path = find('gcia_app/images/gc-h-logo.png')
    title_slide.shapes.add_picture(
        image_path,
        left,
        top,
        width=prs.slide_width,
        height=prs.slide_height
    )
    
    # Add summary slide with overall portfolio information
    add_portfolio_summary_slide(prs, client_details, summary_data, logo_path, transaction_data)
    
    # Add Equity MF Review slide
    add_equity_review_slide(prs, summary_data, logo_path)

    # Add Hybrid MF Review slide
    add_hybrid_review_slide(prs, summary_data, logo_path)

    # Add Debt MF Review slide
    add_debt_review_slide(prs, summary_data, logo_path)
    
    # Get all unique keys for columns
    keys = ['S.No'] + list(ppt_data_dict[0].keys())
    
    # Calculate how many funds can fit on one slide
    funds_per_slide = 4
    
    # Calculate rows needed for each fund
    max_list_length = max(
        max(len(value) if isinstance(value, list) else 1 
            for value in fund.values())
        for fund in ppt_data_dict
    )
    
    # Define three colors for fund rows
    fund_row_colors = [
        RGBColor(220, 230, 242),  # Base blue
        RGBColor(235, 240, 247),  # Lighter blue
        RGBColor(245, 248, 252)   # Lightest blue
    ]
    
    # Split data into chunks for multiple slides
    for slide_start in range(0, len(ppt_data_dict), funds_per_slide):
        slide_data = ppt_data_dict[slide_start:slide_start + funds_per_slide]
        
        # Create content slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add slide title with page numbers if multiple slides
        slide_title = "Mutual Fund Portfolio Review"
        if len(ppt_data_dict) > funds_per_slide:
            current_page = (slide_start // funds_per_slide) + 1
            total_pages = (len(ppt_data_dict) + funds_per_slide - 1) // funds_per_slide
            if total_pages > 1:
                slide_title += f" ({current_page}/{total_pages})"
        
        add_custom_header(
            slide, 
            slide_title, 
            "This Slide Reviews the Mutual Fund Portfolio", 
            logo_path
        )
        
        # Calculate rows needed for this slide
        rows = (max_list_length * len(slide_data)) + 1  # +1 for header
        cols = len(keys)
        
        # Adjusted table dimensions
        left = Inches(0.5)
        top = Inches(1.75)
        width = Inches(13.433)
        height = Inches(6.5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths based on content type
        col_widths = {
            'S.No': 0.4,
            'Scheme Name': 1.8,
            'Purchase Value': 0.9,
            'Current Value': 0.9,
            'Gain': 0.9,
            'Weight': 0.8,
            'Stocks': 0.8,
            'YTD': 0.9,
            '1YR': 0.9,
            '3YR': 0.9,
            '5YR': 0.9,
            'PE': 0.8,
            'Index PE': 0.8,
            'Comments': 1.733
        }
        
        # Add headers and set column widths
        for col_idx, key in enumerate(keys):
            cell = table.cell(0, col_idx)
            cell.text = key
            
            # Header formatting with center alignment
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Set text color to white
            run = paragraph.runs[0]
            run.font.color.rgb = RGBColor(255, 255, 255)
            
            # Set column width
            table.columns[col_idx].width = Inches(col_widths.get(key, 1.0))
            
            # Change header color to royal blue (match other slides)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(65, 105, 225)
        
        
        # Add data
        for fund_idx, fund_data in enumerate(slide_data):
            start_row = 1 + (fund_idx * max_list_length)
            
            # Initialize current_fund_height before processing data
            current_fund_height = max(
                len(value) if isinstance(value, list) else 1 
                for value in fund_data.values()
            )
            
            # Set background colors for all rows in this fund
            for row_idx in range(current_fund_height):
                current_row = start_row + row_idx
                color = RGBColor(220, 230, 242) if fund_idx % 2 == 0 else RGBColor(235, 240, 247)
                
                for col_idx in range(cols):
                    cell = table.cell(current_row, col_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = color
            
            # Add S.No
            cell = table.cell(start_row, 0)
            cell.text = str(slide_start + fund_idx + 1)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Process each column
            for col_idx, key in enumerate(keys[1:], start=1):
                value = fund_data[key]
                
                if isinstance(value, list):
                    for row_idx, item in enumerate(value):
                        cell = table.cell(start_row + row_idx, col_idx)
                        cell.text = str(item)
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(10)
                        paragraph.alignment = PP_ALIGN.CENTER
                        
                        if row_idx < len(value) - 1:
                            paragraph.line_spacing = 1.2
                
                else:
                    cell = table.cell(start_row, col_idx)
                    cell.text = str(value)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.CENTER
                    
                    # Merge cells vertically if needed
                    if current_fund_height > 1:
                        start_cell = cell
                        end_cell = table.cell(start_row + current_fund_height - 1, col_idx)
                        start_cell.merge(end_cell)
                        start_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Merge S.No cells for the current fund
            if current_fund_height > 1:
                start_cell = table.cell(start_row, 0)
                end_cell = table.cell(start_row + current_fund_height - 1, 0)
                start_cell.merge(end_cell)
                start_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Adjust row heights
        for row in table.rows:
            row.height = Inches(0.35)

    # slide = prs.slides.add_slide(prs.slide_layouts[5])
    # add_custom_header(
    #     slide, 
    #     "Mutual Fund Portfolio Review Summary", 
    #     "", 
    #     logo_path,
    #     summary_slide=True
    # )

    # Save the presentation
    prs.save(client_details["Client Name"]+'_fund_analysis.pptx')
    return client_details["Client Name"]+'_fund_analysis.pptx'



